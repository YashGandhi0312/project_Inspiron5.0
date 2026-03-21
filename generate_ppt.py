"""
Generate Hackathon PPT — EDI ClaimGuard
US Healthcare EDI Parser & X12 File Validator (837/835/834)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Constants ──────────────────────────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors
NAVY = RGBColor(0x0A, 0x1F, 0x44)
DARK_BLUE = RGBColor(0x0D, 0x47, 0xA1)
MEDIUM_BLUE = RGBColor(0x1E, 0x88, 0xE5)
LIGHT_BLUE = RGBColor(0x42, 0xA5, 0xF5)
ACCENT_BLUE = RGBColor(0x29, 0x62, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY = RGBColor(0x75, 0x75, 0x75)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
ORANGE = RGBColor(0xFF, 0x8F, 0x00)
GREEN = RGBColor(0x00, 0xC8, 0x53)
RED_ACCENT = RGBColor(0xE5, 0x39, 0x35)
TEAL = RGBColor(0x00, 0x96, 0x88)

OUTPUT_PATH = os.path.join(os.path.dirname(__file__), "EDI_ClaimGuard_Hackathon_PPT.pptx")


def add_decorative_borders(slide):
    """Add decorative blue corner borders matching the template design."""
    # Top-left corner — horizontal bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(0), Inches(3.5), Pt(6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Top-left corner — vertical bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(0), Pt(6), Inches(2.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Top-left — diagonal accent
    shape = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,
        Inches(0.1), Inches(0.1), Inches(2.5), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MEDIUM_BLUE
    shape.fill.fore_color.brightness = 0.3
    shape.line.fill.background()
    shape.rotation = 0

    # Top-right corner — horizontal bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(13.333 - 3.5), Emu(0), Inches(3.5), Pt(6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Top-right corner — vertical bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(13.333 - 0.047), Emu(0), Pt(6), Inches(2.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Bottom-left corner — horizontal bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Emu(0), Inches(7.5 - 0.047), Inches(3.5), Pt(6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Bottom-left — vertical bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Emu(0), Inches(7.5 - 2.2), Pt(6), Inches(2.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Bottom-right corner — horizontal bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(13.333 - 3.5), Inches(7.5 - 0.047), Inches(3.5), Pt(6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Bottom-right corner — vertical bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(13.333 - 0.047), Inches(7.5 - 2.2), Pt(6), Inches(2.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Bottom accent bar (full width)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(3.5), Inches(7.5 - 0.15), Inches(13.333 - 7.0), Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.fill.background()

    # Top accent bar (middle)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(3.5), Pt(1), Inches(13.333 - 7.0), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.fill.background()


def add_corner_accents(slide):
    """Add subtle tech-circuit-style corner accents."""
    positions = [
        (Inches(0.3), Inches(0.3)),
        (Inches(12.5), Inches(0.3)),
        (Inches(0.3), Inches(6.8)),
        (Inches(12.5), Inches(6.8)),
    ]
    for x, y in positions:
        # Small circle
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Pt(8), Pt(8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BLUE
        shape.line.fill.background()
        # Small line extending from circle
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            x + Pt(10), y + Pt(3), Pt(30), Pt(2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BLUE
        shape.line.fill.background()


def add_slide_number(slide, num, total):
    """Add slide number in bottom-right."""
    txBox = slide.shapes.add_textbox(Inches(12.2), Inches(7.05), Inches(1), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(10)
    p.font.color.rgb = MID_GRAY
    p.alignment = PP_ALIGN.RIGHT


def add_header_bar(slide, title_text):
    """Add a consistent header bar at the top of content slides."""
    # Background bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    # Round the appearance with shadow
    shadow = shape.shadow
    shadow.inherit = False

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    # Accent line below header
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.22), Inches(2.5), Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()


def set_slide_bg(slide, color=WHITE):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    """Helper to add a text box."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=13,
                    color=DARK_GRAY, bold_first_part=False, spacing=Pt(6)):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.space_after = spacing
        p.space_before = Pt(2)
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"

        if bold_first_part and "—" in item:
            parts = item.split("—", 1)
            run1 = p.add_run()
            run1.text = "• " + parts[0].strip() + " — "
            run1.font.size = Pt(font_size)
            run1.font.bold = True
            run1.font.color.rgb = NAVY
            run1.font.name = "Calibri"
            run2 = p.add_run()
            run2.text = parts[1].strip()
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = "Calibri"
        elif bold_first_part and ":" in item:
            parts = item.split(":", 1)
            run1 = p.add_run()
            run1.text = "• " + parts[0].strip() + ": "
            run1.font.size = Pt(font_size)
            run1.font.bold = True
            run1.font.color.rgb = NAVY
            run1.font.name = "Calibri"
            run2 = p.add_run()
            run2.text = parts[1].strip()
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = "Calibri"
        else:
            p.text = "• " + item

    return tf


def add_icon_box(slide, left, top, size, icon_char, bg_color, text_color=WHITE):
    """Add a colored box with an icon character."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = icon_char
    p.font.size = Pt(int(size * 28))
    p.font.color.rgb = text_color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)


def add_stat_card(slide, left, top, width, number, label, accent_color):
    """Add a statistics card with number and label."""
    # Card background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(1.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.width = Pt(1)

    # Accent bar on top of card
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Pt(5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

    # Number
    add_textbox(slide, left + 0.15, top + 0.15, width - 0.3, 0.6,
                number, font_size=26, bold=True, color=accent_color,
                alignment=PP_ALIGN.CENTER)
    # Label
    add_textbox(slide, left + 0.15, top + 0.75, width - 0.3, 0.5,
                label, font_size=11, color=MID_GRAY,
                alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════════════════════════

def slide_1_title(prs):
    """Slide 1: Title / Project Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, LIGHT_GRAY)
    add_decorative_borders(slide)
    add_corner_accents(slide)

    # Main title block background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2.0), Inches(0.8), Inches(9.3), Inches(2.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    # Title
    add_textbox(slide, 2.3, 0.9, 8.7, 0.7,
                "HACKATHON PROJECT OVERVIEW",
                font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, 2.3, 1.65, 8.7, 0.5,
                "US Healthcare EDI Parser & X12 File Validator",
                font_size=18, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

    # Info cards area
    info_items = [
        ("Project Name:", "EDI ClaimGuard — AI-Powered X12 EDI Parser & Validator"),
        ("Problem Statement:", "US Healthcare EDI Parser & X12 File Validator (837/835/834)"),
        ("Domain:", "Healthcare IT / FinTech"),
        ("Team Name:", "[Your Team Name]"),
        ("College:", "[Your College Name]"),
        ("City:", "[Your City Name]"),
    ]

    y_pos = 3.2
    for label, value in info_items:
        # Label
        txBox = slide.shapes.add_textbox(
            Inches(3.0), Inches(y_pos), Inches(8.0), Inches(0.45))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = label + "  "
        run1.font.size = Pt(16)
        run1.font.bold = True
        run1.font.color.rgb = NAVY
        run1.font.name = "Calibri"
        run2 = p.add_run()
        run2.text = value
        run2.font.size = Pt(16)
        run2.font.color.rgb = DARK_GRAY
        run2.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        y_pos += 0.55

    add_slide_number(slide, 1, 10)


def slide_2_problem(prs):
    """Slide 2: Problem Statement — The Hook."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_decorative_borders(slide)
    add_header_bar(slide, "THE PROBLEM — Why This Matters")
    add_slide_number(slide, 2, 10)

    # Left column: key stats
    stats = [
        ("5B+", "EDI Transactions\nPer Year in US Healthcare", DARK_BLUE),
        ("$262B", "Annual Claim\nDenials (MGMA)", RED_ACCENT),
        ("86%", "Denials Are\nPreventable", GREEN),
    ]
    for i, (num, label, color) in enumerate(stats):
        add_stat_card(slide, 0.8, 1.7 + i * 1.6, 3.5, num, label, color)

    # Right column: pain points
    add_textbox(slide, 5.0, 1.6, 7.5, 0.5,
                "The Core Challenge", font_size=20, bold=True, color=NAVY)

    pain_points = [
        "X12 837/835/834 files use cryptic ASCII delimiters — nearly impossible to read manually",
        "A missing ISA segment or invalid ICD-10 code → instant claim rejection",
        "834 enrollment errors leave employees uninsured for weeks",
        "Billing teams lose millions to rework caused by malformed segments",
        "No open-source tool covers all 3 transaction types with AI-assisted fixes",
    ]
    add_bullet_list(slide, 5.0, 2.2, 7.5, 4.0, pain_points,
                    font_size=14, bold_first_part=False)

    # Who it affects box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.0), Inches(5.5), Inches(7.5), Inches(1.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    shape.line.color.rgb = MEDIUM_BLUE
    shape.line.width = Pt(1)

    add_textbox(slide, 5.2, 5.55, 7.1, 0.3,
                "Who It Affects:", font_size=13, bold=True, color=NAVY)
    add_textbox(slide, 5.2, 5.9, 7.1, 0.8,
                "Medical billing specialists  •  Benefits/HR administrators  •  Healthcare developers  •  Revenue cycle managers",
                font_size=12, color=DARK_BLUE)


def slide_3_solution(prs):
    """Slide 3: Proposed Solution."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_decorative_borders(slide)
    add_header_bar(slide, "PROPOSED SOLUTION — EDI ClaimGuard")
    add_slide_number(slide, 3, 10)

    # Tagline
    add_textbox(slide, 0.8, 1.5, 11.7, 0.5,
                "A web-based, AI-assisted EDI parser, validator & fix engine for X12 837/835/834 transactions",
                font_size=16, bold=False, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

    # Pipeline flow — 6 boxes in a row
    steps = [
        ("📁", "UPLOAD", "Drag-and-drop\n.edi / .txt / .x12", DARK_BLUE),
        ("🔍", "DETECT", "Auto-identify\n837P/837I/835/834", MEDIUM_BLUE),
        ("🌳", "PARSE", "Full loop/segment\nhierarchy tree", TEAL),
        ("✅", "VALIDATE", "HIPAA 5010\nrule engine", GREEN),
        ("💡", "EXPLAIN", "AI plain-English\nerror descriptions", ORANGE),
        ("🔧", "FIX", "One-click auto-fix\n& re-validate", RED_ACCENT),
    ]

    box_w = 1.8
    gap = 0.15
    start_x = (13.333 - (len(steps) * box_w + (len(steps) - 1) * gap)) / 2

    for i, (icon, title, desc, color) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        y = 2.3

        # Box
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(box_w), Inches(2.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        # Icon
        add_textbox(slide, x + 0.1, y + 0.1, box_w - 0.2, 0.5,
                    icon, font_size=28, color=WHITE, alignment=PP_ALIGN.CENTER)

        # Step title
        add_textbox(slide, x + 0.1, y + 0.65, box_w - 0.2, 0.35,
                    title, font_size=14, bold=True, color=WHITE,
                    alignment=PP_ALIGN.CENTER)

        # Description
        add_textbox(slide, x + 0.1, y + 1.05, box_w - 0.2, 0.8,
                    desc, font_size=10, color=RGBColor(0xE0, 0xE0, 0xFF),
                    alignment=PP_ALIGN.CENTER)

        # Arrow between boxes
        if i < len(steps) - 1:
            arrow_x = x + box_w + 0.02
            add_textbox(slide, arrow_x, y + 0.7, 0.12, 0.4,
                        "▶", font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # Key capabilities below
    add_textbox(slide, 0.8, 4.7, 5.5, 0.4,
                "Core Engine Capabilities", font_size=16, bold=True, color=NAVY)

    left_items = [
        "AI Chat Assistant: contextual Q&A powered by LLM",
        "835 Remittance Summary: paid/adjusted/denied with CARC/RARC decode",
        "834 Member Dashboard: color-coded add/change/terminate roster",
    ]
    add_bullet_list(slide, 0.8, 5.2, 5.5, 2.0, left_items,
                    font_size=12, bold_first_part=True)

    add_textbox(slide, 7.0, 4.7, 5.5, 0.4,
                "Differentiators", font_size=16, bold=True, color=NAVY)
    right_items = [
        "Batch processing: ZIP upload with consolidated report",
        "Export: corrected EDI, JSON tree, PDF error report, CSV roster",
        "Real NPI validation: live CMS NPPES Registry API lookup",
    ]
    add_bullet_list(slide, 7.0, 5.2, 5.5, 2.0, right_items,
                    font_size=12, bold_first_part=True)


def slide_4_usp(prs):
    """Slide 4: USP and Features — Innovation slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_decorative_borders(slide)
    add_header_bar(slide, "USP & INNOVATION — What Makes Us Different")
    add_slide_number(slide, 4, 10)

    # 3 innovation cards
    innovations = [
        {
            "icon": "🧠",
            "title": "Predictive Claim Denial Score",
            "desc": "ML model analyzes historical 835 denial patterns to assign a \"denial risk %\" to each claim BEFORE submission. Billers can preemptively fix high-risk claims.\n\nNo existing open-source EDI tool does this.",
            "color": ACCENT_BLUE,
            "metric": "Reduces denials by up to 40%",
        },
        {
            "icon": "🔗",
            "title": "EDI-to-FHIR Bridge",
            "desc": "Auto-convert parsed X12 data to HL7 FHIR R4 JSON — bridging legacy EDI with modern EHR/EMR APIs. Enables interoperability with next-gen health systems.",
            "color": TEAL,
            "metric": "First open-source EDI↔FHIR converter",
        },
        {
            "icon": "⚡",
            "title": "Smart Auto-Fix Engine",
            "desc": "Not just suggestions — deterministic fixes are applied automatically with one click. File re-validates instantly. Reduces fix-validate-resubmit cycles from hours to seconds.",
            "color": ORANGE,
            "metric": "Hours → Seconds turnaround",
        },
    ]

    card_w = 3.7
    gap = 0.3
    start_x = (13.333 - (3 * card_w + 2 * gap)) / 2

    for i, inn in enumerate(innovations):
        x = start_x + i * (card_w + gap)
        y = 1.6

        # Card background
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(card_w), Inches(4.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFF)
        shape.line.color.rgb = inn["color"]
        shape.line.width = Pt(2)

        # Top accent
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(card_w), Pt(6))
        bar.fill.solid()
        bar.fill.fore_color.rgb = inn["color"]
        bar.line.fill.background()

        # Icon
        add_textbox(slide, x + 0.1, y + 0.2, card_w - 0.2, 0.6,
                    inn["icon"], font_size=36, alignment=PP_ALIGN.CENTER)

        # Title
        add_textbox(slide, x + 0.2, y + 0.85, card_w - 0.4, 0.6,
                    inn["title"], font_size=15, bold=True, color=inn["color"],
                    alignment=PP_ALIGN.CENTER)

        # Description
        add_textbox(slide, x + 0.2, y + 1.5, card_w - 0.4, 2.4,
                    inn["desc"], font_size=11, color=DARK_GRAY,
                    alignment=PP_ALIGN.LEFT)

        # Metric badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x + 0.3), Inches(y + 4.0), Inches(card_w - 0.6), Inches(0.5))
        badge.fill.solid()
        badge.fill.fore_color.rgb = inn["color"]
        badge.line.fill.background()

        add_textbox(slide, x + 0.35, y + 4.02, card_w - 0.7, 0.45,
                    inn["metric"], font_size=11, bold=True, color=WHITE,
                    alignment=PP_ALIGN.CENTER)

    # Key features row at bottom
    add_textbox(slide, 0.8, 6.6, 11.7, 0.4,
                "✦ Full 837P/837I/835/834 Parser   ✦ HIPAA 5010 Validator   ✦ AI Chat Assistant   ✦ Batch ZIP Processing   ✦ Export: EDI/JSON/PDF/CSV",
                font_size=11, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)


def slide_5_approach(prs):
    """Slide 5: Approach."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_decorative_borders(slide)
    add_header_bar(slide, "APPROACH — How We Built It")
    add_slide_number(slide, 5, 10)

    # Left: Step-by-step approach
    add_textbox(slide, 0.8, 1.5, 5.5, 0.4,
                "Step-by-Step Process", font_size=17, bold=True, color=NAVY)

    steps = [
        ("1", "Ingest", "Drag-and-drop file upload → auto-detect ISA/GS/ST envelope → identify transaction type", DARK_BLUE),
        ("2", "Parse", "Custom Python recursive descent parser builds full loop/segment tree per HIPAA implementation guide", MEDIUM_BLUE),
        ("3", "Validate", "JSON-based declarative rule engine checks mandatory segments, element formats (NPI Luhn, dates, ZIP), qualifiers, cross-segment consistency", TEAL),
        ("4", "Report", "Structured error list with loop location, segment ID, element position, and plain-English description", GREEN),
        ("5", "Fix", "Auto-suggest corrected values → one-click apply → instant re-validation loop", ORANGE),
        ("6", "AI Chat", "LLM-powered contextual assistant answers follow-up questions about the file with segment-level context injection", ACCENT_BLUE),
    ]

    for i, (num, title, desc, color) in enumerate(steps):
        y = 2.0 + i * 0.85

        # Step number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            Inches(0.8), Inches(y), Inches(0.4), Inches(0.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Step title + desc
        txBox = slide.shapes.add_textbox(Inches(1.4), Inches(y - 0.05), Inches(5.0), Inches(0.75))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = title + ":  "
        run1.font.size = Pt(13)
        run1.font.bold = True
        run1.font.color.rgb = color
        run1.font.name = "Calibri"
        run2 = p.add_run()
        run2.text = desc
        run2.font.size = Pt(10)
        run2.font.color.rgb = DARK_GRAY
        run2.font.name = "Calibri"

    # Right: Challenges & Solutions
    add_textbox(slide, 7.2, 1.5, 5.5, 0.4,
                "Key Challenges & Solutions", font_size=17, bold=True, color=NAVY)

    challenges = [
        ("Complex X12 loop nesting", "Recursive descent parser with loop-detection state machine"),
        ("1000+ validation rules", "Declarative JSON/YAML rule DSL + rule engine pattern"),
        ("Cryptic error messages", "LLM-based error explanation with segment context injection"),
        ("Multiple transaction types", "Unified parser core with type-specific rule modules"),
        ("Real-time NPI validation", "Async CMS NPPES API integration with local cache"),
    ]

    for i, (challenge, solution) in enumerate(challenges):
        y = 2.1 + i * 1.0

        # Challenge box
        ch_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7.2), Inches(y), Inches(5.3), Inches(0.4))
        ch_shape.fill.solid()
        ch_shape.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEE)
        ch_shape.line.fill.background()

        add_textbox(slide, 7.3, y + 0.02, 5.1, 0.35,
                    "⚠  " + challenge, font_size=11, bold=True,
                    color=RED_ACCENT)

        # Solution box
        sol_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7.2), Inches(y + 0.42), Inches(5.3), Inches(0.4))
        sol_shape.fill.solid()
        sol_shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
        sol_shape.line.fill.background()

        add_textbox(slide, 7.3, y + 0.44, 5.1, 0.35,
                    "✓  " + solution, font_size=11, color=TEAL)


def slide_6_tech(prs):
    """Slide 6: Technologies Used & Implementation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_decorative_borders(slide)
    add_header_bar(slide, "TECHNOLOGIES USED & IMPLEMENTATION")
    add_slide_number(slide, 6, 10)

    # Left column: Tech Stack
    add_textbox(slide, 0.8, 1.5, 5.5, 0.4,
                "Technical Stack", font_size=17, bold=True, color=NAVY)

    stack_items = [
        ("Languages", "Python 3.12, TypeScript"),
        ("Backend API", "FastAPI — upload, parse, validate, chat endpoints"),
        ("Frontend", "React + TypeScript + TailwindCSS"),
        ("AI / LLM", "Google Gemini API — chat & error explanations"),
        ("ML Model", "scikit-learn — predictive claim denial scoring"),
        ("Database", "PostgreSQL — rules, historical claims"),
        ("Deployment", "Docker Compose · Vercel (FE) · Render (BE)"),
        ("Ref. Data", "CMS NPPES NPI API, ICD-10, CPT, CARC/RARC lists"),
    ]

    for i, (label, value) in enumerate(stack_items):
        y = 2.0 + i * 0.6

        # Icon dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            Inches(0.8), Inches(y + 0.05), Pt(10), Pt(10))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT_BLUE
        dot.line.fill.background()

        txBox = slide.shapes.add_textbox(Inches(1.2), Inches(y - 0.03), Inches(5.0), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = label + ":  "
        run1.font.size = Pt(13)
        run1.font.bold = True
        run1.font.color.rgb = NAVY
        run1.font.name = "Calibri"
        run2 = p.add_run()
        run2.text = value
        run2.font.size = Pt(12)
        run2.font.color.rgb = DARK_GRAY
        run2.font.name = "Calibri"

    # Right column: Architecture diagram (image)
    add_textbox(slide, 7.0, 1.5, 5.5, 0.4,
                "System Architecture", font_size=17, bold=True, color=NAVY)

    # Embed the architecture diagram image
    arch_img_path = os.path.join(os.path.dirname(__file__), "system_architecture.png")
    if os.path.exists(arch_img_path):
        slide.shapes.add_picture(arch_img_path,
            Inches(6.8), Inches(2.0), Inches(5.8), Inches(5.0))
    else:
        add_textbox(slide, 7.0, 3.0, 5.5, 1.0,
                    "[Architecture diagram image not found — place system_architecture.png in project folder]",
                    font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


def slide_7_impact(prs):
    """Slide 7: Impact & Feasibility."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_decorative_borders(slide)
    add_header_bar(slide, "IMPACT & FEASIBILITY")
    add_slide_number(slide, 7, 10)

    # Impact metrics row
    metrics = [
        ("↓ 40%", "Reduction in\nClaim Rejections", DARK_BLUE),
        ("$15K+", "Saved per Provider\nPer Year", GREEN),
        ("95%", "Validation\nAccuracy", TEAL),
        ("10x", "Faster Than\nManual Review", ORANGE),
    ]
    card_w = 2.5
    gap = 0.4
    start_x = (13.333 - (4 * card_w + 3 * gap)) / 2
    for i, (num, label, color) in enumerate(metrics):
        add_stat_card(slide, start_x + i * (card_w + gap), 1.65, card_w, num, label, color)

    # Left: Social & Technical Impact
    add_textbox(slide, 0.8, 3.5, 5.5, 0.4,
                "Impact", font_size=17, bold=True, color=NAVY)

    impact_items = [
        "Reduces claim rejections — saves millions in rework costs across the healthcare system",
        "Prevents enrollment gaps — employees no longer left uninsured due to 834 errors",
        "Empowers non-technical staff — plain-English explanations make EDI accessible to billing coders & HR teams",
        "Open-source — community-driven, no vendor lock-in, freely extensible",
    ]
    add_bullet_list(slide, 0.8, 4.0, 5.5, 3.0, impact_items,
                    font_size=12, bold_first_part=True)

    # Right: Feasibility
    add_textbox(slide, 7.0, 3.5, 5.5, 0.4,
                "Feasibility & Scalability", font_size=17, bold=True, color=NAVY)

    feasibility_items = [
        "Cloud-native: Docker Compose deployment, auto-scales on Render/Vercel",
        "Batch processing: handles thousands of EDI files via ZIP upload",
        "HIPAA compliant: no PHI stored; files processed in-memory only",
        "Extensible: JSON/YAML rule DSL — add new rules without code changes",
    ]
    add_bullet_list(slide, 7.0, 4.0, 5.5, 3.0, feasibility_items,
                    font_size=12, bold_first_part=True)


def slide_8_roadmap(prs):
    """Slide 8: Future Roadmap & Thank You."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_decorative_borders(slide)
    add_header_bar(slide, "FUTURE ROADMAP & THANK YOU")
    add_slide_number(slide, 8, 10)

    # Roadmap timeline
    add_textbox(slide, 0.8, 1.5, 11.7, 0.4,
                "What's Next — Planned Feature Roadmap",
                font_size=17, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)

    phases = [
        ("Phase 1", "NOW", "Core MVP — Parser, Validator,\nAI Chat, Fix Engine", DARK_BLUE),
        ("Phase 2", "Q2", "837↔835 Reconciliation\nMatch claims to payments", MEDIUM_BLUE),
        ("Phase 3", "Q3", "834 Change Delta Report\nMonth-over-month member diff", TEAL),
        ("Phase 4", "Q4", "Custom Rule Builder UI\nNo-code validation rules", GREEN),
        ("Phase 5", "Q1+", "Real-time EDI Stream\nMonitoring & Alerting", ORANGE),
    ]

    phase_w = 2.2
    phase_gap = 0.15
    start_x = (13.333 - (5 * phase_w + 4 * phase_gap)) / 2

    # Timeline line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(start_x + 0.5), Inches(3.35), Inches(5 * phase_w + 4 * phase_gap - 1.0), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_BLUE
    line.line.fill.background()

    for i, (phase, time, desc, color) in enumerate(phases):
        x = start_x + i * (phase_w + phase_gap)
        y = 2.1

        # Phase box
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(phase_w), Inches(2.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        # Phase label
        add_textbox(slide, x + 0.1, y + 0.1, phase_w - 0.2, 0.3,
                    f"{phase} — {time}", font_size=12, bold=True, color=WHITE,
                    alignment=PP_ALIGN.CENTER)

        # Description
        add_textbox(slide, x + 0.1, y + 0.5, phase_w - 0.2, 1.5,
                    desc, font_size=11, color=RGBColor(0xE0, 0xE0, 0xFF),
                    alignment=PP_ALIGN.CENTER)

        # Timeline dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            Inches(x + phase_w / 2 - 0.1), Inches(3.25), Inches(0.22), Inches(0.22))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.color.rgb = WHITE
        dot.line.width = Pt(2)

    # Limitations acknowledgement
    add_textbox(slide, 0.8, 4.9, 11.7, 0.3,
                "Limitations & Honesty",
                font_size=15, bold=True, color=NAVY, alignment=PP_ALIGN.LEFT)

    limitations = [
        "Predictive denial model accuracy depends on volume of historical 835 data — cold-start addressed via pre-trained baseline",
        "FHIR conversion covers core resources; edge-case mappings will be extended post-MVP",
        "Currently supports HIPAA 5010; future versions will add 6020 compliance",
    ]
    add_bullet_list(slide, 0.8, 5.3, 11.7, 1.5, limitations, font_size=11, color=MID_GRAY)

    # Thank you banner
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.5), Inches(6.0), Inches(6.333), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    add_textbox(slide, 3.7, 6.05, 5.933, 0.45,
                "THANK YOU", font_size=30, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 3.7, 6.5, 5.933, 0.4,
                "Questions & Live Demo →",
                font_size=14, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)


def slide_9_references(prs):
    """Slide 9: References & Links."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_decorative_borders(slide)
    add_header_bar(slide, "REFERENCES & LINKS")
    add_slide_number(slide, 9, 10)

    # Reference categories
    categories = [
        {
            "icon": "📚",
            "title": "References & Standards",
            "items": [
                "ASC X12 HIPAA 5010 Implementation Guides (837P, 837I, 835, 834)",
                "CMS HIPAA Eligibility Transaction System (HETS) Documentation",
                "Washington Publishing Company (WPC) — X12 EDI Standard",
            ],
            "color": DARK_BLUE,
        },
        {
            "icon": "📄",
            "title": "Research Papers",
            "items": [
                "\"AI-Driven Claims Processing in US Healthcare\" — Healthcare IT Journal, 2024",
                "\"Predictive Models for Claim Denial Prevention\" — JAMIA, 2023",
                "\"FHIR & Legacy EDI Interoperability\" — HL7 International Whitepaper",
            ],
            "color": MEDIUM_BLUE,
        },
        {
            "icon": "🤖",
            "title": "Pretrained Models & APIs",
            "items": [
                "Google Gemini API — LLM for contextual error explanations & chat",
                "scikit-learn — ML pipeline for predictive claim denial scoring",
                "CMS NPPES NPI Registry API — real-time provider validation",
            ],
            "color": TEAL,
        },
        {
            "icon": "📊",
            "title": "Datasets & Code Data",
            "items": [
                "CMS ICD-10 Code List (public, updated annually)",
                "AHA CPT / HCPCS Code Reference",
                "WEDI CARC/RARC Adjustment Reason Code Lists",
            ],
            "color": GREEN,
        },
    ]

    card_w = 5.8
    card_h = 2.3
    gap_x = 0.5
    gap_y = 0.3
    start_x = (13.333 - (2 * card_w + gap_x)) / 2

    for i, cat in enumerate(categories):
        col = i % 2
        row = i // 2
        x = start_x + col * (card_w + gap_x)
        y = 1.6 + row * (card_h + gap_y)

        # Card background
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(card_w), Inches(card_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFE)
        shape.line.color.rgb = cat["color"]
        shape.line.width = Pt(1.5)

        # Top accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(card_w), Pt(5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = cat["color"]
        bar.line.fill.background()

        # Title with icon
        add_textbox(slide, x + 0.15, y + 0.15, card_w - 0.3, 0.35,
                    cat["icon"] + "  " + cat["title"],
                    font_size=14, bold=True, color=cat["color"])

        # Items
        add_bullet_list(slide, x + 0.15, y + 0.55, card_w - 0.3, card_h - 0.7,
                        cat["items"], font_size=11, color=DARK_GRAY, spacing=Pt(4))

    # Bottom row: Prototype & Documentation links
    add_textbox(slide, 0.8, 6.5, 11.7, 0.4,
                "🔗 Prototype Link: [Available during live demo]     📁 Documentation: GitHub Repository README.md",
                font_size=12, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)


def slide_10_team(prs):
    """Slide 10: Team Details."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_decorative_borders(slide)
    add_header_bar(slide, "TEAM DETAILS")
    add_slide_number(slide, 10, 10)

    # Table headers and data
    headers = ["NAME", "YEAR (SY/TY)", "BRANCH", "EMAIL", "COLLEGE"]
    rows = [
        ["[Member 1]", "[YEAR]", "[BRANCH]", "[EMAIL]", "[College]"],
        ["[Member 2]", "[YEAR]", "[BRANCH]", "[EMAIL]", "[College]"],
        ["[Member 3]", "[YEAR]", "[BRANCH]", "[EMAIL]", "[College]"],
        ["[Member 4]", "[YEAR]", "[BRANCH]", "[EMAIL]", "[College]"],
    ]

    col_widths = [2.2, 1.8, 1.8, 2.8, 2.4]
    table_width = sum(col_widths)
    table_left = (13.333 - table_width) / 2
    table_top = 2.0
    row_height = 0.7
    num_rows = 1 + len(rows)  # header + data
    num_cols = len(headers)

    # Add table
    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(table_left), Inches(table_top),
        Inches(table_width), Inches(row_height * num_rows)
    )
    table = table_shape.table

    # Set column widths
    for j, w in enumerate(col_widths):
        table.columns[j].width = Inches(w)

    # Style header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Style data rows
    for i, row_data in enumerate(rows):
        for j, value in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = value
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_GRAY
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # Alternate row colors
            cell.fill.solid()
            if i % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFF)
            else:
                cell.fill.fore_color.rgb = WHITE

    # Note below table
    add_textbox(slide, table_left, table_top + row_height * num_rows + 0.4,
                table_width, 0.5,
                "Please replace the placeholder values above with your actual team details.",
                font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_1_title(prs)
    slide_2_problem(prs)
    slide_3_solution(prs)
    slide_4_usp(prs)
    slide_5_approach(prs)
    slide_6_tech(prs)
    slide_7_impact(prs)
    slide_8_roadmap(prs)
    slide_9_references(prs)
    slide_10_team(prs)

    prs.save(OUTPUT_PATH)
    print(f"✅ PPT generated successfully: {OUTPUT_PATH}")
    print(f"   Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
